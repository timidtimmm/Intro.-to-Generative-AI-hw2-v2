"""
tools.py — Extended Tool Registry (MCP-style)
All existing tools preserved + new MCP tools added.
"""

import base64
import hashlib
import json
import math
import os
import re
import time
import uuid
from dataclasses import dataclass
from datetime import datetime, timezone, timedelta
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

import httpx

# ── Data model ────────────────────────────────────────────────────────────────

@dataclass
class ToolSpec:
    name: str
    description: str
    parameters: Dict
    handler: Callable
    risk_level: str = "low"
    timeout_s: int = 15
    require_approval: bool = False
    category: str = "general"
    icon: str = "🔧"


# ── Registry ──────────────────────────────────────────────────────────────────

class ToolRegistry:
    def __init__(self):
        self._tools: Dict[str, ToolSpec] = {}

    def register(self, spec: ToolSpec):
        self._tools[spec.name] = spec

    def get(self, name: str) -> Optional[ToolSpec]:
        return self._tools.get(name)

    def list_tools(self) -> List[Dict]:
        return [
            {
                "name": s.name,
                "description": s.description,
                "risk_level": s.risk_level,
                "require_approval": s.require_approval,
                "category": s.category,
                "icon": s.icon,
            }
            for s in self._tools.values()
        ]

    def openai_tool_defs(self, names: Optional[List[str]] = None) -> List[Dict]:
        specs = (
            [self._tools[n] for n in names if n in self._tools]
            if names else list(self._tools.values())
        )
        return [
            {
                "type": "function",
                "function": {
                    "name": s.name,
                    "description": s.description,
                    "parameters": s.parameters,
                },
            }
            for s in specs
        ]

    def execute(self, name: str, args: Dict[str, Any]) -> Dict[str, Any]:
        spec = self.get(name)
        if not spec:
            return {"error": f"Unknown tool: {name}"}
        t0 = time.time()
        try:
            result = spec.handler(args)
        except Exception as exc:
            result = {"error": str(exc)}
        result["_latency_ms"] = round((time.time() - t0) * 1000, 1)
        return result


# ══════════════════════════════════════════════════════════════════════════════
#  TOOL IMPLEMENTATIONS
# ══════════════════════════════════════════════════════════════════════════════

# ── 1. Web Search ─────────────────────────────────────────────────────────────
def _web_search(args: Dict) -> Dict:
    query = args.get("query", "").strip()
    if not query:
        return {"error": "query is required"}
    try:
        from duckduckgo_search import DDGS
        with DDGS() as ddgs:
            hits = list(ddgs.text(query, max_results=5))
        return {
            "results": [
                {"title": h.get("title",""), "url": h.get("href",""), "snippet": h.get("body","")[:400]}
                for h in hits
            ],
            "query": query,
        }
    except ImportError:
        return {"error": "pip install duckduckgo-search"}
    except Exception as e:
        return {"error": str(e), "query": query}


# ── 2. GitHub Repo Search ─────────────────────────────────────────────────────
def _github_repo_search(args: Dict) -> Dict:
    query = args.get("query", "").strip()
    per_page = min(max(int(args.get("per_page", 5) or 5), 1), 10)
    if not query:
        return {"error": "query is required"}
    headers = {"Accept": "application/vnd.github+json", "User-Agent": "hw2-chatgpt-demo"}
    token = os.getenv("GITHUB_TOKEN", "").strip()
    if token:
        headers["Authorization"] = f"Bearer {token}"
    try:
        resp = httpx.get(
            "https://api.github.com/search/repositories",
            params={"q": query, "sort": "stars", "order": "desc", "per_page": per_page},
            headers=headers, timeout=20.0,
        )
        data = resp.json()
        if resp.status_code >= 400:
            return {"error": data.get("message", resp.text[:500]), "query": query}
        items = [
            {"full_name": i.get("full_name"), "description": i.get("description"),
             "html_url": i.get("html_url"), "stars": i.get("stargazers_count"),
             "language": i.get("language"), "updated_at": i.get("updated_at")}
            for i in data.get("items", [])[:per_page]
        ]
        return {"query": query, "count": len(items), "results": items}
    except Exception as e:
        return {"error": str(e), "query": query}


# ── 3. Calculator ─────────────────────────────────────────────────────────────
def _calculator(args: Dict) -> Dict:
    expr = args.get("expression", "").strip()
    if not expr:
        return {"error": "expression is required"}
    safe = re.sub(r"[^0-9\+\-\*\/\^\(\)\.\s]", "", expr).replace("^", "**")
    if not safe.strip():
        return {"error": "Invalid expression"}
    try:
        allowed = {k: v for k, v in math.__dict__.items() if not k.startswith("__")}
        result = eval(safe, {"__builtins__": {}}, allowed)  # nosec
        return {"result": result, "expression": expr, "formatted": f"{expr} = {result}"}
    except ZeroDivisionError:
        return {"error": "除以零"}
    except Exception as e:
        return {"error": str(e)}


# ── 4. Weather ────────────────────────────────────────────────────────────────
def _get_weather(args: Dict) -> Dict:
    location = args.get("location", "").strip()
    if not location:
        return {"error": "location is required"}
    try:
        resp = httpx.get(f"https://wttr.in/{location}?format=j1", timeout=10.0, headers={"User-Agent": "curl/7.64.1"})
        data = resp.json()
        cur = data.get("current_condition", [{}])[0]
        area = data.get("nearest_area", [{}])[0]
        area_name = area.get("areaName", [{}])[0].get("value", location)
        country = area.get("country", [{}])[0].get("value", "")
        return {
            "location": f"{area_name}, {country}",
            "description": cur.get("weatherDesc", [{}])[0].get("value", "?"),
            "temperature_c": cur.get("temp_C", "?"),
            "feels_like_c": cur.get("FeelsLikeC", "?"),
            "humidity_pct": cur.get("humidity", "?"),
            "wind_kmph": cur.get("windspeedKmph", "?"),
        }
    except Exception as e:
        return {"error": str(e), "location": location}


# ── 5. Datetime ───────────────────────────────────────────────────────────────
def _get_datetime(args: Dict) -> Dict:
    from datetime import datetime as _datetime
    try:
        from zoneinfo import ZoneInfo
    except ImportError:
        try:
            from backports.zoneinfo import ZoneInfo  # type: ignore
        except ImportError:
            utc_plus8 = timezone(timedelta(hours=8))
            now = _datetime.now(utc_plus8)
            return {
                "timezone": "Asia/Taipei (UTC+8 fallback)", "datetime": now.strftime("%Y-%m-%d %H:%M:%S"),
                "date": now.strftime("%Y-%m-%d"), "time": now.strftime("%H:%M:%S"),
                "weekday": now.strftime("%A"),
                "weekday_zh": ["星期一","星期二","星期三","星期四","星期五","星期六","星期日"][now.weekday()],
                "unix_ts": int(now.timestamp()),
            }
    tz_name = (args.get("timezone") or "Asia/Taipei").strip()
    try:
        tz = ZoneInfo(tz_name)
    except Exception:
        tz = ZoneInfo("Asia/Taipei"); tz_name = "Asia/Taipei (fallback)"
    now = _datetime.now(tz)
    return {
        "timezone": tz_name, "datetime": now.strftime("%Y-%m-%d %H:%M:%S"),
        "date": now.strftime("%Y-%m-%d"), "time": now.strftime("%H:%M:%S"),
        "weekday": now.strftime("%A"),
        "weekday_zh": ["星期一","星期二","星期三","星期四","星期五","星期六","星期日"][now.weekday()],
        "unix_ts": int(now.timestamp()),
    }


# ── 6. Image Generation ───────────────────────────────────────────────────────
def _gemini_api_key() -> str:
    return os.getenv("GEMINI_API_KEY", "").strip()

_IMAGE_MODEL_CANDIDATES = [
    "gemini-2.0-flash-exp",
    "gemini-2.0-flash-exp-image-generation",
    "gemini-2.0-flash",
]

def _call_gemini_generate_content(model: str, payload: Dict) -> Dict:
    api_key = _gemini_api_key()
    if not api_key:
        return {"error": "GEMINI_API_KEY not set"}
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    try:
        resp = httpx.post(url, json=payload, timeout=120.0)
        data = resp.json()
        if resp.status_code >= 400:
            return {"error": data.get("error", {}).get("message", resp.text[:1000]), "model": model}
        return data
    except Exception as e:
        return {"error": str(e), "model": model}

def _extract_gemini_inline_image(data: Dict) -> Optional[Dict]:
    for cand in data.get("candidates", []) or []:
        for part in (cand.get("content", {}) or {}).get("parts", []) or []:
            inline = part.get("inlineData") or part.get("inline_data")
            if isinstance(inline, dict) and inline.get("data"):
                return {"mime": inline.get("mimeType") or "image/png", "image_b64": inline["data"]}
    return None

def _extract_gemini_text(data: Dict) -> str:
    texts = []
    for cand in data.get("candidates", []) or []:
        for part in (cand.get("content", {}) or {}).get("parts", []) or []:
            if isinstance(part, dict) and part.get("text"):
                texts.append(part["text"])
    return "\n".join(t.strip() for t in texts if t).strip()

def _generate_image(args: Dict) -> Dict:
    prompt = args.get("prompt", "").strip()
    if not prompt:
        return {"error": "prompt is required"}
    api_key = _gemini_api_key()
    if not api_key:
        return {"error": "GEMINI_API_KEY 未設定"}
    env_override = os.getenv("GEMINI_IMAGE_MODEL", "").strip()
    candidates = ([env_override] if env_override else []) + [m for m in _IMAGE_MODEL_CANDIDATES if m != env_override]
    payload = {"contents": [{"role": "user", "parts": [{"text": prompt}]}], "generationConfig": {"responseModalities": ["TEXT", "IMAGE"]}}
    last_error = "No model available"
    for model in candidates:
        data = _call_gemini_generate_content(model, payload)
        if "error" not in data:
            img = _extract_gemini_inline_image(data)
            if img:
                result = {**img, "prompt": prompt, "model": model}
                text = _extract_gemini_text(data)
                if text:
                    result["caption"] = text
                return result
            return {"error": f"Model {model} returned no image. Reply: {_extract_gemini_text(data)[:200]}", "prompt": prompt}
        last_error = str(data.get("error", "unknown"))
        lower_err = last_error.lower()
        if any(k in lower_err for k in ["api key", "quota", "billing", "unauthorized"]):
            return {"error": f"Gemini API 錯誤: {last_error}", "prompt": prompt}
    return {"error": f"All image models failed. Last: {last_error}", "prompt": prompt}


# ── 7. Create Presentation ────────────────────────────────────────────────────
def _create_presentation(args: Dict) -> Dict:
    title = (args.get("title") or "New Presentation").strip()
    slides = args.get("slides") or []
    subtitle = (args.get("subtitle") or "").strip()
    filename = (args.get("filename") or "").strip()
    if not isinstance(slides, list) or not slides:
        return {"error": "slides must be a non-empty array"}
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:
        return {"error": "pip install python-pptx"}
    _base = Path(__file__).resolve().parent
    out_dir = _base / "static" / "generated"
    out_dir.mkdir(parents=True, exist_ok=True)
    safe_name = re.sub(r"[^a-zA-Z0-9_\-\u4e00-\u9fff]+", "_", filename or title).strip("_") or f"ppt_{uuid.uuid4().hex[:8]}"
    out_path = out_dir / f"{safe_name}.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle or "Generated by MCP Tools"
    for item in slides:
        if not isinstance(item, dict):
            continue
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = str(item.get("title") or "")
        body = s.shapes.placeholders[1].text_frame
        body.clear()
        added = False
        for bullet in (item.get("bullets") or []):
            p = body.paragraphs[0] if not added else body.add_paragraph()
            p.text = str(bullet).strip(); p.level = 0; p.font.size = Pt(20); added = True
        if not added and item.get("notes"):
            body.text = str(item["notes"])
        if item.get("notes"):
            notes_slide = s.notes_slide
            notes_slide.notes_text_frame.text = str(item["notes"])
    prs.save(str(out_path))
    return {"title": title, "slide_count": len(prs.slides), "file_path": str(out_path), "download_url": f"/api/download/{out_path.name}"}


# ── 8. URL Fetch ──────────────────────────────────────────────────────────────
def _url_fetch(args: Dict) -> Dict:
    url = (args.get("url") or "").strip()
    if not url:
        return {"error": "url is required"}
    if not url.startswith(("http://", "https://")):
        url = "https://" + url
    max_chars = min(int(args.get("max_chars", 3000)), 10000)
    try:
        resp = httpx.get(url, timeout=15.0, follow_redirects=True,
                         headers={"User-Agent": "Mozilla/5.0 (compatible; hw2-bot/1.0)"})
        content_type = resp.headers.get("content-type", "")
        if "text/html" in content_type:
            # Strip HTML tags
            text = re.sub(r"<style[^>]*>[\s\S]*?</style>", "", resp.text, flags=re.I)
            text = re.sub(r"<script[^>]*>[\s\S]*?</script>", "", text, flags=re.I)
            text = re.sub(r"<[^>]+>", " ", text)
            text = re.sub(r"\s+", " ", text).strip()
        elif "json" in content_type:
            text = json.dumps(resp.json(), ensure_ascii=False, indent=2)
        else:
            text = resp.text
        return {
            "url": str(resp.url),
            "status_code": resp.status_code,
            "content_type": content_type,
            "content": text[:max_chars],
            "truncated": len(text) > max_chars,
            "content_length": len(text),
        }
    except Exception as e:
        return {"error": str(e), "url": url}


# ── 9. Unit Converter ─────────────────────────────────────────────────────────
_UNITS: Dict[str, Dict[str, float]] = {
    # Length → meters
    "m": 1, "km": 1000, "cm": 0.01, "mm": 0.001, "mile": 1609.344,
    "yd": 0.9144, "ft": 0.3048, "in": 0.0254, "nm": 1852,
    # Weight → grams
    "g": 1, "kg": 1000, "mg": 0.001, "lb": 453.592, "oz": 28.3495, "ton": 1e6,
    # Temperature (special)
    "c": None, "f": None, "k": None,
    # Speed → m/s
    "ms": 1, "kmh": 1/3.6, "mph": 0.44704, "knot": 0.514444,
    # Area → m²
    "m2": 1, "km2": 1e6, "cm2": 1e-4, "ha": 1e4, "acre": 4046.856,
    # Volume → liters
    "l": 1, "ml": 0.001, "gal": 3.78541, "qt": 0.946353, "cup": 0.236588, "tbsp": 0.0147868, "tsp": 0.00492892,
    # Digital → bytes
    "b": 1, "kb": 1024, "mb": 1024**2, "gb": 1024**3, "tb": 1024**4,
}
_UNIT_GROUPS = {
    "length": {"m","km","cm","mm","mile","yd","ft","in","nm"},
    "weight": {"g","kg","mg","lb","oz","ton"},
    "temperature": {"c","f","k"},
    "speed": {"ms","kmh","mph","knot"},
    "area": {"m2","km2","cm2","ha","acre"},
    "volume": {"l","ml","gal","qt","cup","tbsp","tsp"},
    "digital": {"b","kb","mb","gb","tb"},
}

def _unit_converter(args: Dict) -> Dict:
    value = args.get("value")
    from_unit = str(args.get("from_unit", "")).lower().strip()
    to_unit = str(args.get("to_unit", "")).lower().strip()
    if value is None or not from_unit or not to_unit:
        return {"error": "value, from_unit, to_unit are required"}
    try:
        value = float(value)
    except Exception:
        return {"error": "value must be a number"}

    # Temperature conversion
    if from_unit in {"c","f","k"} or to_unit in {"c","f","k"}:
        def to_celsius(v, u):
            if u == "c": return v
            if u == "f": return (v - 32) * 5/9
            if u == "k": return v - 273.15
        def from_celsius(v, u):
            if u == "c": return v
            if u == "f": return v * 9/5 + 32
            if u == "k": return v + 273.15
        result = from_celsius(to_celsius(value, from_unit), to_unit)
        return {"result": round(result, 6), "from": f"{value} {from_unit.upper()}", "to": f"{round(result,4)} {to_unit.upper()}"}

    if from_unit not in _UNITS or to_unit not in _UNITS:
        return {"error": f"Unknown unit(s): {from_unit}, {to_unit}. Available: {', '.join(sorted(_UNITS.keys()))}"}

    # Check same group
    from_group = next((g for g, us in _UNIT_GROUPS.items() if from_unit in us), None)
    to_group = next((g for g, us in _UNIT_GROUPS.items() if to_unit in us), None)
    if from_group != to_group:
        return {"error": f"Cannot convert between different types: {from_group} → {to_group}"}

    base = value * _UNITS[from_unit]
    result = base / _UNITS[to_unit]
    return {
        "result": round(result, 8),
        "from": f"{value} {from_unit}",
        "to": f"{round(result, 6)} {to_unit}",
        "category": from_group,
        "formula": f"{value} × {_UNITS[from_unit]} / {_UNITS[to_unit]}",
    }


# ── 10. Note Manager ──────────────────────────────────────────────────────────
_NOTES: Dict[str, Dict] = {}  # In-memory; persists for server lifetime

def _note_manager(args: Dict) -> Dict:
    action = (args.get("action") or "list").lower().strip()
    if action == "list":
        notes = sorted(_NOTES.values(), key=lambda n: n["updated_at"], reverse=True)
        return {"notes": notes, "count": len(notes)}
    if action == "create":
        title = (args.get("title") or "Untitled").strip()
        content = (args.get("content") or "").strip()
        nid = uuid.uuid4().hex[:8]
        now = datetime.utcnow().isoformat()
        _NOTES[nid] = {"id": nid, "title": title, "content": content, "created_at": now, "updated_at": now}
        return {"created": True, "note": _NOTES[nid]}
    if action == "read":
        nid = args.get("id", "").strip()
        if nid not in _NOTES:
            # Search by title
            nid = next((k for k, v in _NOTES.items() if args.get("title","").lower() in v["title"].lower()), None)
            if not nid:
                return {"error": "Note not found"}
        return {"note": _NOTES[nid]}
    if action == "update":
        nid = args.get("id", "").strip()
        if nid not in _NOTES:
            return {"error": "Note not found"}
        if args.get("title"):
            _NOTES[nid]["title"] = args["title"]
        if args.get("content") is not None:
            _NOTES[nid]["content"] = args["content"]
        _NOTES[nid]["updated_at"] = datetime.utcnow().isoformat()
        return {"updated": True, "note": _NOTES[nid]}
    if action == "delete":
        nid = args.get("id", "").strip()
        if nid in _NOTES:
            note = _NOTES.pop(nid)
            return {"deleted": True, "note": note}
        return {"error": "Note not found"}
    if action == "search":
        q = (args.get("query") or "").lower()
        results = [n for n in _NOTES.values() if q in n["title"].lower() or q in n["content"].lower()]
        return {"results": results, "count": len(results)}
    return {"error": f"Unknown action: {action}. Use: list, create, read, update, delete, search"}


# ── 11. QR Code Generator ─────────────────────────────────────────────────────
def _qr_generator(args: Dict) -> Dict:
    text = (args.get("text") or "").strip()
    if not text:
        return {"error": "text is required"}
    size = min(max(int(args.get("size", 200)), 50), 1000)
    # Use goqr.me API (free, no key needed)
    try:
        encoded = httpx.URL(text)  # validate URL
        url = f"https://api.qrserver.com/v1/create-qr-code/?size={size}x{size}&data={httpx.URL(text)}"
        resp = httpx.get(url, timeout=15.0)
        if resp.status_code == 200:
            b64 = base64.b64encode(resp.content).decode()
            return {
                "qr_image_b64": b64,
                "mime": "image/png",
                "text": text,
                "size": size,
                "api_url": str(resp.url),
            }
        return {"error": f"QR API returned {resp.status_code}", "text": text}
    except Exception as e:
        return {"error": str(e), "text": text}


# ── 12. Timezone Converter ────────────────────────────────────────────────────
def _timezone_converter(args: Dict) -> Dict:
    from datetime import datetime as _dt
    from_tz = (args.get("from_timezone") or "Asia/Taipei").strip()
    to_tz = (args.get("to_timezone") or "UTC").strip()
    dt_str = args.get("datetime") or ""
    try:
        from zoneinfo import ZoneInfo
    except ImportError:
        return {"error": "zoneinfo not available (Python 3.9+ required)"}
    try:
        from_zone = ZoneInfo(from_tz)
        to_zone = ZoneInfo(to_tz)
    except Exception as e:
        return {"error": f"Invalid timezone: {e}"}
    try:
        if dt_str:
            # Try parsing common formats
            for fmt in ["%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M", "%Y-%m-%dT%H:%M:%S", "%Y-%m-%d"]:
                try:
                    dt = _dt.strptime(dt_str, fmt).replace(tzinfo=from_zone)
                    break
                except ValueError:
                    continue
            else:
                return {"error": f"Cannot parse datetime: {dt_str}. Use YYYY-MM-DD HH:MM:SS"}
        else:
            dt = _dt.now(from_zone)
        converted = dt.astimezone(to_zone)
        offset_diff = converted.utcoffset().total_seconds() - dt.utcoffset().total_seconds()
        return {
            "from": {"timezone": from_tz, "datetime": dt.strftime("%Y-%m-%d %H:%M:%S %Z"), "weekday": dt.strftime("%A")},
            "to": {"timezone": to_tz, "datetime": converted.strftime("%Y-%m-%d %H:%M:%S %Z"), "weekday": converted.strftime("%A")},
            "offset_diff_hours": offset_diff / 3600,
        }
    except Exception as e:
        return {"error": str(e)}


# ── 13. Hash Generator ────────────────────────────────────────────────────────
def _hash_generator(args: Dict) -> Dict:
    text = args.get("text") or args.get("input", "")
    if not text:
        return {"error": "text is required"}
    algos = str(args.get("algorithms", "all")).lower()
    text_bytes = text.encode("utf-8")
    if algos == "all":
        return {
            "input": text[:200],
            "md5": hashlib.md5(text_bytes).hexdigest(),
            "sha1": hashlib.sha1(text_bytes).hexdigest(),
            "sha256": hashlib.sha256(text_bytes).hexdigest(),
            "sha512": hashlib.sha512(text_bytes).hexdigest(),
        }
    result = {"input": text[:200]}
    for algo in algos.split(","):
        algo = algo.strip()
        try:
            result[algo] = hashlib.new(algo, text_bytes).hexdigest()
        except ValueError:
            result[algo] = f"unsupported algorithm: {algo}"
    return result


# ── 14. Color Picker / Converter ──────────────────────────────────────────────
def _color_converter(args: Dict) -> Dict:
    color = str(args.get("color") or "").strip()
    if not color:
        return {"error": "color is required"}

    def hex_to_rgb(h):
        h = h.lstrip("#")
        if len(h) == 3:
            h = "".join(c*2 for c in h)
        if len(h) != 6:
            raise ValueError(f"Invalid hex: #{h}")
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    def rgb_to_hsl(r, g, b):
        r, g, b = r/255, g/255, b/255
        mx, mn = max(r,g,b), min(r,g,b)
        l = (mx + mn) / 2
        if mx == mn:
            h = s = 0.0
        else:
            d = mx - mn
            s = d / (2 - mx - mn) if l > 0.5 else d / (mx + mn)
            if mx == r: h = (g - b) / d + (6 if g < b else 0)
            elif mx == g: h = (b - r) / d + 2
            else: h = (r - g) / d + 4
            h /= 6
        return round(h*360), round(s*100), round(l*100)

    try:
        if color.startswith("#"):
            r, g, b = hex_to_rgb(color)
        elif color.lower().startswith("rgb"):
            nums = re.findall(r"\d+", color)
            r, g, b = int(nums[0]), int(nums[1]), int(nums[2])
        else:
            # Named colors
            named = {"red":(255,0,0),"green":(0,128,0),"blue":(0,0,255),"white":(255,255,255),
                     "black":(0,0,0),"yellow":(255,255,0),"cyan":(0,255,255),"magenta":(255,0,255),
                     "orange":(255,165,0),"purple":(128,0,128),"pink":(255,192,203)}
            if color.lower() in named:
                r, g, b = named[color.lower()]
            else:
                return {"error": f"Cannot parse color: {color}. Use #RRGGBB, rgb(r,g,b), or named colors"}

        h, s, l = rgb_to_hsl(r, g, b)
        lum = 0.2126 * r/255 + 0.7152 * g/255 + 0.0722 * b/255
        return {
            "input": color,
            "hex": f"#{r:02x}{g:02x}{b:02x}",
            "HEX": f"#{r:02X}{g:02X}{b:02X}",
            "rgb": {"r": r, "g": g, "b": b, "css": f"rgb({r}, {g}, {b})"},
            "hsl": {"h": h, "s": s, "l": l, "css": f"hsl({h}, {s}%, {l}%)"},
            "luminance": round(lum, 4),
            "is_dark": lum < 0.5,
            "complementary": f"#{(255-r):02x}{(255-g):02x}{(255-b):02x}",
        }
    except Exception as e:
        return {"error": str(e), "input": color}


# ── 15. Dictionary / Translation ──────────────────────────────────────────────
def _dictionary_lookup(args: Dict) -> Dict:
    word = (args.get("word") or "").strip()
    if not word:
        return {"error": "word is required"}
    lang = (args.get("lang") or "en").strip()
    try:
        if lang == "en":
            resp = httpx.get(f"https://api.dictionaryapi.dev/api/v2/entries/en/{word}", timeout=10.0)
            if resp.status_code == 404:
                return {"error": f"Word '{word}' not found in dictionary"}
            data = resp.json()
            if not data or not isinstance(data, list):
                return {"error": "No data returned"}
            entry = data[0]
            result = {
                "word": entry.get("word", word),
                "phonetic": entry.get("phonetic", ""),
                "meanings": [],
                "synonyms": [],
                "antonyms": [],
            }
            for meaning in entry.get("meanings", [])[:3]:
                pos = meaning.get("partOfSpeech", "")
                defs = [d.get("definition","") for d in meaning.get("definitions", [])[:3]]
                result["meanings"].append({"part_of_speech": pos, "definitions": defs})
                result["synonyms"].extend(meaning.get("synonyms", [])[:5])
                result["antonyms"].extend(meaning.get("antonyms", [])[:5])
            result["synonyms"] = list(set(result["synonyms"]))[:10]
            result["antonyms"] = list(set(result["antonyms"]))[:10]
            return result
        else:
            return {"error": f"Language '{lang}' not yet supported. Only 'en' currently available."}
    except Exception as e:
        return {"error": str(e), "word": word}


# ── 16. Random Generator ──────────────────────────────────────────────────────
def _random_generator(args: Dict) -> Dict:
    import random, string
    kind = (args.get("type") or "number").lower().strip()
    if kind == "number":
        lo = int(args.get("min", 1))
        hi = int(args.get("max", 100))
        count = min(int(args.get("count", 1)), 100)
        nums = [random.randint(lo, hi) for _ in range(count)]
        return {"type": "number", "results": nums if count > 1 else nums[0], "range": f"{lo}–{hi}", "count": count}
    if kind == "password":
        length = min(max(int(args.get("length", 16)), 4), 128)
        chars = string.ascii_letters + string.digits + "!@#$%^&*()-_=+"
        return {"type": "password", "result": "".join(random.choices(chars, k=length)), "length": length}
    if kind == "uuid":
        count = min(int(args.get("count", 1)), 20)
        uuids = [str(uuid.uuid4()) for _ in range(count)]
        return {"type": "uuid", "results": uuids if count > 1 else uuids[0]}
    if kind == "name":
        first = ["Alice","Bob","Carol","David","Emma","Frank","Grace","Henry","Iris","Jack","Kate","Leo","Mia","Nina","Oscar","Petra","Quinn","Rosa","Sam","Tina","Ugo","Vera","Will","Xena","Yuki","Zoe"]
        last = ["Smith","Jones","Brown","Wilson","Davis","Miller","Moore","Taylor","Anderson","Thomas","White","Harris","Martin","Garcia","Lee","Robinson","Walker","Hall","Allen","Young"]
        count = min(int(args.get("count", 1)), 10)
        names = [f"{random.choice(first)} {random.choice(last)}" for _ in range(count)]
        return {"type": "name", "results": names if count > 1 else names[0]}
    if kind == "color":
        count = min(int(args.get("count", 1)), 20)
        colors = [f"#{random.randint(0,0xFFFFFF):06x}" for _ in range(count)]
        return {"type": "color", "results": colors if count > 1 else colors[0]}
    if kind == "choice":
        choices = args.get("choices", [])
        if isinstance(choices, str):
            choices = [c.strip() for c in choices.split(",")]
        if not choices:
            return {"error": "choices is required for type=choice"}
        count = min(int(args.get("count", 1)), len(choices))
        picked = random.sample(choices, count) if count > 1 else random.choice(choices)
        return {"type": "choice", "results": picked, "from": choices}
    return {"error": f"Unknown type: {kind}. Use: number, password, uuid, name, color, choice"}


# ══════════════════════════════════════════════════════════════════════════════
#  REGISTRY SETUP
# ══════════════════════════════════════════════════════════════════════════════

registry = ToolRegistry()

registry.register(ToolSpec(
    name="web_search", icon="🔍", category="information",
    description="Search the internet for current information, recent news, or facts.",
    parameters={"type":"object","properties":{"query":{"type":"string","description":"Search query"}},"required":["query"]},
    handler=_web_search, risk_level="low", timeout_s=15,
))

registry.register(ToolSpec(
    name="github_repo_search", icon="🐙", category="developer",
    description="Search public GitHub repositories for APIs, SDKs, open-source examples.",
    parameters={"type":"object","properties":{"query":{"type":"string"},"per_page":{"type":"integer","default":5}},"required":["query"]},
    handler=_github_repo_search, risk_level="low", timeout_s=20,
))

registry.register(ToolSpec(
    name="calculator", icon="🧮", category="math",
    description="Evaluate math expressions: arithmetic, algebra, unit conversions.",
    parameters={"type":"object","properties":{"expression":{"type":"string","description":"e.g. '2 * (3 + 4) / 7'"}},"required":["expression"]},
    handler=_calculator, risk_level="low", timeout_s=5,
))

registry.register(ToolSpec(
    name="get_weather", icon="🌤", category="information",
    description="Get current weather for any city.",
    parameters={"type":"object","properties":{"location":{"type":"string","description":"City name e.g. 'Taipei'"}},"required":["location"]},
    handler=_get_weather, risk_level="low", timeout_s=10,
))

registry.register(ToolSpec(
    name="get_datetime", icon="🕐", category="utility",
    description="Get current date/time for any timezone.",
    parameters={"type":"object","properties":{"timezone":{"type":"string","description":"IANA timezone e.g. 'Asia/Taipei'"}},"required":[]},
    handler=_get_datetime, risk_level="low", timeout_s=5,
))

registry.register(ToolSpec(
    name="generate_image", icon="🎨", category="media",
    description="Generate an image from text using Gemini AI.",
    parameters={"type":"object","properties":{"prompt":{"type":"string","description":"Detailed image description"}},"required":["prompt"]},
    handler=_generate_image, risk_level="medium", timeout_s=120,
))

registry.register(ToolSpec(
    name="create_presentation", icon="📊", category="productivity",
    description="Create a PowerPoint (.pptx) presentation file.",
    parameters={"type":"object","properties":{
        "title":{"type":"string"},"subtitle":{"type":"string"},"filename":{"type":"string"},
        "slides":{"type":"array","items":{"type":"object","properties":{"title":{"type":"string"},"bullets":{"type":"array","items":{"type":"string"}},"notes":{"type":"string"}},"required":["title","bullets"]}},
    },"required":["title","slides"]},
    handler=_create_presentation, risk_level="medium", timeout_s=30,
))

registry.register(ToolSpec(
    name="url_fetch", icon="🌐", category="information",
    description="Fetch and read content from a URL (web page, JSON API, etc.).",
    parameters={"type":"object","properties":{
        "url":{"type":"string","description":"URL to fetch"},
        "max_chars":{"type":"integer","description":"Max characters to return (default 3000)","default":3000},
    },"required":["url"]},
    handler=_url_fetch, risk_level="low", timeout_s=20,
))

registry.register(ToolSpec(
    name="unit_converter", icon="📐", category="utility",
    description="Convert units: length, weight, temperature, speed, area, volume, digital storage.",
    parameters={"type":"object","properties":{
        "value":{"type":"number","description":"Value to convert"},
        "from_unit":{"type":"string","description":"Source unit e.g. 'km', 'lb', 'f'"},
        "to_unit":{"type":"string","description":"Target unit e.g. 'm', 'kg', 'c'"},
    },"required":["value","from_unit","to_unit"]},
    handler=_unit_converter, risk_level="low", timeout_s=5,
))

registry.register(ToolSpec(
    name="note_manager", icon="📝", category="productivity",
    description="Create, read, update, delete, and search notes. Actions: list, create, read, update, delete, search.",
    parameters={"type":"object","properties":{
        "action":{"type":"string","enum":["list","create","read","update","delete","search"]},
        "id":{"type":"string"}, "title":{"type":"string"}, "content":{"type":"string"}, "query":{"type":"string"},
    },"required":["action"]},
    handler=_note_manager, risk_level="low", timeout_s=5,
))

registry.register(ToolSpec(
    name="qr_generator", icon="📱", category="utility",
    description="Generate a QR code image for any text or URL.",
    parameters={"type":"object","properties":{
        "text":{"type":"string","description":"Text or URL to encode"},
        "size":{"type":"integer","description":"Image size in pixels (50-1000, default 200)","default":200},
    },"required":["text"]},
    handler=_qr_generator, risk_level="low", timeout_s=15,
))

registry.register(ToolSpec(
    name="timezone_converter", icon="🌍", category="utility",
    description="Convert date/time between different timezones.",
    parameters={"type":"object","properties":{
        "from_timezone":{"type":"string","description":"Source IANA timezone e.g. 'Asia/Taipei'"},
        "to_timezone":{"type":"string","description":"Target IANA timezone e.g. 'America/New_York'"},
        "datetime":{"type":"string","description":"Optional: datetime string YYYY-MM-DD HH:MM:SS (defaults to now)"},
    },"required":["from_timezone","to_timezone"]},
    handler=_timezone_converter, risk_level="low", timeout_s=5,
))

registry.register(ToolSpec(
    name="hash_generator", icon="🔒", category="developer",
    description="Generate MD5, SHA1, SHA256, SHA512 hashes for any text.",
    parameters={"type":"object","properties":{
        "text":{"type":"string","description":"Text to hash"},
        "algorithms":{"type":"string","description":"Comma-separated or 'all' (default: all)","default":"all"},
    },"required":["text"]},
    handler=_hash_generator, risk_level="low", timeout_s=5,
))

registry.register(ToolSpec(
    name="color_converter", icon="🎨", category="developer",
    description="Convert colors between HEX, RGB, HSL formats and get complementary color.",
    parameters={"type":"object","properties":{
        "color":{"type":"string","description":"Color in #RRGGBB, rgb(r,g,b), or named format"},
    },"required":["color"]},
    handler=_color_converter, risk_level="low", timeout_s=5,
))

registry.register(ToolSpec(
    name="dictionary_lookup", icon="📖", category="language",
    description="Look up English word definitions, phonetics, synonyms, and antonyms.",
    parameters={"type":"object","properties":{
        "word":{"type":"string","description":"Word to look up"},
        "lang":{"type":"string","description":"Language (currently only 'en')","default":"en"},
    },"required":["word"]},
    handler=_dictionary_lookup, risk_level="low", timeout_s=10,
))

registry.register(ToolSpec(
    name="random_generator", icon="🎲", category="utility",
    description="Generate random numbers, passwords, UUIDs, names, colors, or choices.",
    parameters={"type":"object","properties":{
        "type":{"type":"string","enum":["number","password","uuid","name","color","choice"],"description":"What to generate"},
        "min":{"type":"integer","description":"Min number (for type=number)"},
        "max":{"type":"integer","description":"Max number (for type=number)"},
        "length":{"type":"integer","description":"Password length (for type=password)"},
        "count":{"type":"integer","description":"How many to generate"},
        "choices":{"description":"Array or CSV string of choices (for type=choice)"},
    },"required":["type"]},
    handler=_random_generator, risk_level="low", timeout_s=5,
))

import re as _re
 
def _extract_video_id(url: str) -> str:
    """從各種 YouTube URL 格式抓出 11 碼 video ID。"""
    patterns = [
        r"(?:v=|youtu\.be/)([A-Za-z0-9_-]{11})",
        r"(?:embed/)([A-Za-z0-9_-]{11})",
        r"(?:shorts/)([A-Za-z0-9_-]{11})",
    ]
    for p in patterns:
        m = _re.search(p, url)
        if m:
            return m.group(1)
    # 若傳入的本身就是 11 碼 ID
    if _re.match(r"^[A-Za-z0-9_-]{11}$", url.strip()):
        return url.strip()
    return ""
 
 
def _youtube_transcript(args: Dict) -> Dict:
    """
    取得 YouTube 影片字幕並選擇性摘要。
    
    args:
      url      (str) YouTube 網址或影片 ID
      lang     (str) 優先語言，預設 zh-Hant,zh-Hans,en
      summarize(bool) 是否回傳摘要（字數較少），預設 False
      max_chars(int) 回傳字幕最大字元數，預設 3000
    """
    url = (args.get("url") or args.get("video_id") or "").strip()
    if not url:
        return {"error": "url 或 video_id 為必填"}
 
    video_id = _extract_video_id(url)
    if not video_id:
        return {"error": f"無法從 '{url}' 解析出 YouTube 影片 ID"}
 
    lang_pref = (args.get("lang") or "zh-Hant,zh-Hans,en").split(",")
    lang_pref = [l.strip() for l in lang_pref if l.strip()]
    max_chars = min(int(args.get("max_chars", 3000)), 10000)
    do_summary = str(args.get("summarize", "false")).lower() in ("true", "1", "yes")
 
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
    except ImportError:
        return {"error": "請先安裝：pip install youtube-transcript-api --break-system-packages"}
 
    # ── 嘗試各語言 ────────────────────────────────────────────────────────────
    transcript_obj = None
    used_lang = None
    ytt = YouTubeTranscriptApi()
 
    try:
        transcript_list = ytt.list(video_id)
    except Exception as e:
        return {
            "error": f"無法取得字幕清單：{e}",
            "video_id": video_id,
            "youtube_url": f"https://www.youtube.com/watch?v={video_id}",
        }
 
    # 依優先順序找字幕
    for lang in lang_pref:
        try:
            transcript_obj = transcript_list.find_transcript([lang])
            used_lang = lang
            break
        except Exception:
            continue
 
    # fallback：自動生成字幕
    if not transcript_obj:
        try:
            transcript_obj = transcript_list.find_generated_transcript(lang_pref)
            used_lang = "auto-generated"
        except Exception:
            pass
 
    # 最後 fallback：任何語言
    if not transcript_obj:
        try:
            all_transcripts = list(transcript_list)
            if all_transcripts:
                transcript_obj = all_transcripts[0]
                used_lang = transcript_obj.language_code
        except Exception as e:
            return {
                "error": f"找不到任何字幕：{e}",
                "video_id": video_id,
                "youtube_url": f"https://www.youtube.com/watch?v={video_id}",
            }
 
    if not transcript_obj:
        return {
            "error": "此影片沒有任何字幕（可能是私人影片或字幕未開放）",
            "video_id": video_id,
        }
 
    # ── 取得字幕文字 ──────────────────────────────────────────────────────────
    try:
        fetched = transcript_obj.fetch()
        snippets = fetched.snippets  # 新版 API
    except AttributeError:
        # 舊版 API fallback
        fetched = transcript_obj.fetch()
        snippets = fetched
 
    full_text = " ".join(
        (s.text if hasattr(s, "text") else s.get("text", ""))
        for s in snippets
    ).strip()
 
    # ── 摘要模式：只回傳前 max_chars 字，並標注時間戳 ────────────────────────
    if do_summary:
        # 每 30 秒取一個代表句子
        timeline = []
        step = 30
        bucket: list = []
        current_bucket_start = 0
 
        for s in snippets:
            start = s.start if hasattr(s, "start") else s.get("start", 0)
            text  = s.text  if hasattr(s, "text")  else s.get("text", "")
            if start >= current_bucket_start + step:
                if bucket:
                    timeline.append(f"[{int(current_bucket_start//60):02d}:{int(current_bucket_start%60):02d}] " + " ".join(bucket))
                bucket = [text]
                current_bucket_start = (start // step) * step
            else:
                bucket.append(text)
        if bucket:
            timeline.append(f"[{int(current_bucket_start//60):02d}:{int(current_bucket_start%60):02d}] " + " ".join(bucket))
 
        summary_text = "\n".join(timeline)[:max_chars]
        return {
            "video_id": video_id,
            "youtube_url": f"https://www.youtube.com/watch?v={video_id}",
            "language": used_lang,
            "mode": "timeline_summary",
            "content": summary_text,
            "total_chars": len(full_text),
        }
 
    # ── 一般模式：回傳完整字幕（截斷） ──────────────────────────────────────
    return {
        "video_id": video_id,
        "youtube_url": f"https://www.youtube.com/watch?v={video_id}",
        "language": used_lang,
        "mode": "full_transcript",
        "content": full_text[:max_chars],
        "truncated": len(full_text) > max_chars,
        "total_chars": len(full_text),
    }
 
registry.register(ToolSpec(
    name="youtube_transcript",
    icon="▶️",
    category="media",
    description=(
        "取得 YouTube 影片的字幕文字或時間軸摘要。"
        "支援繁體中文、簡體中文、英文字幕，可自動選擇可用語言。"
        "適合：摘要影片內容、搜尋影片重點、分析演講或教學影片。"
    ),
    parameters={
        "type": "object",
        "properties": {
            "url": {
                "type": "string",
                "description": "YouTube 影片網址或 11 碼影片 ID，例如 https://www.youtube.com/watch?v=xxxxx",
            },
            "lang": {
                "type": "string",
                "description": "優先字幕語言，逗號分隔，預設 'zh-Hant,zh-Hans,en'",
                "default": "zh-Hant,zh-Hans,en",
            },
            "summarize": {
                "type": "boolean",
                "description": "True = 回傳時間軸摘要（每30秒一段），False = 完整字幕",
                "default": False,
            },
            "max_chars": {
                "type": "integer",
                "description": "回傳最大字元數，預設 3000，最大 10000",
                "default": 3000,
            },
        },
        "required": ["url"],
    },
    handler=_youtube_transcript,
    risk_level="low",
    timeout_s=30,
))
